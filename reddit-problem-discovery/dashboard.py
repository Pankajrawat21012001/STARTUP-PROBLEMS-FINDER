"""
Reddit Problem Discovery Dashboard

Run with: streamlit run dashboard.py
"""

import os
import re
import json
import html as html_lib
import pandas as pd
import streamlit as st
import plotly.graph_objects as go
from datetime import datetime
import io
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# ── Page Config ──────────────────────────────────────────────────
st.set_page_config(
    page_title="Problem Discovery Dashboard",
    page_icon="🔍",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ── Paths ────────────────────────────────────────────────────────
DATA_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "data")
PROBLEM_IDS_PATH    = os.path.join(DATA_DIR, "problem_ids.csv")
PROBLEM_EVIDENCE_PATH = os.path.join(DATA_DIR, "problem_evidence.csv")
PROBLEM_SCORES_PATH = os.path.join(DATA_DIR, "problem_scores.csv")
RAW_POSTS_PATH      = os.path.join(DATA_DIR, "raw_posts.csv")
IDEA_EVALUATION_PATH = os.path.join(DATA_DIR, "idea_evaluation.csv")
IDEA_REVIEWS_PATH = os.path.join(DATA_DIR, "idea_reviews.csv")
IDEA_REVIEWS_COLUMNS = ["problem_id", "status", "reason", "notes", "reviewed_at"]

EVALUATION_DIMENSIONS = [
    {"key": "problem_need",      "category": "Problem & Need (Acute Problem)",  "question": "Vitamin or Painkiller? Painful enough to pay for?"},
    {"key": "customer_clarity",  "category": "Customer Clarity",                "question": "Clearly defined target user with WTP?"},
    {"key": "market_size",       "category": "Market Size (Enough?)",           "question": "Market large enough? Growing or shrinking?"},
    {"key": "competition",       "category": "Competition",                     "question": "Who are competitors? Why choose this?"},
    {"key": "demand_validation", "category": "Do People Want This?",            "question": "Organic demand signals exist?"},
    {"key": "recently_possible", "category": "Recently Possible",               "question": "What tech/market shift enables this now?"},
    {"key": "good_proxies",      "category": "Good Proxies",                    "question": "Adjacent companies proving the market?"},
    {"key": "ideaspace",         "category": "Good Ideaspace",                  "question": "Strong category with expansion paths?"},
    {"key": "real_problem",      "category": "Is It a Real Problem?",           "question": "Genuine problem or solution looking for one?"},
    {"key": "tarpit_risk",       "category": "Tarpit Idea",                     "question": "Many failed here before? Why will this succeed?"},
]

# WTP phrases (kept in sync with step4_wtp_score.py)
WTP_PHRASES = [
    "would pay", "i'd pay", "willing to pay", "worth paying",
    "subscription", "enterprise plan", "our company would", "budget for",
    "paying for", "charge for this", "affordable tool", "worth the money",
    "would buy", "need to purchase",
    "take my money", "shut up and take", "where do i pay",
    "how much does it cost", "is there a paid", "premium version",
    "monthly fee", "annual plan", "per month", "pricing",
    "worth it", "definitely pay", "pay for this",
    "invest in", "spend money on", "need this badly",
    "money well spent", "i wish someone would build",
    "why isn't there a", "why is there no",
    "can't believe there's no", "desperately looking for",
    "tried everything", "no good solution", "nothing works",
    "existing tools are terrible", "current solution is broken",
    "sick of manually", "tired of doing this manually",
    "waste hours", "wastes my time", "time consuming",
    "need a tool", "need software", "need an app",
    "is there an app", "any app for this", "any tool for",
]

URGENCY_PHRASES = [
    "urgent", "critical", "costing us", "losing money", "hours wasted",
    "every day", "daily problem", "manually every", "broken process",
    "nightmare", "so frustrated", "desperately need",
    "massive pain", "huge problem", "real pain point",
    "huge time sink", "takes forever", "kills productivity",
    "inefficient", "waste of time", "painful process",
    "no solution", "workaround", "band-aid solution",
    "fed up", "sick of", "really annoying",
    "tedious", "repetitive task", "manual work",
    "this is ridiculous", "why is this so hard",
]


# ── Helper: Load data ────────────────────────────────────────────
@st.cache_data(ttl=60)
def load_data():
    """Load all CSV files with graceful handling for missing files."""
    dfs = {}
    for name, path in [
        ("problems", PROBLEM_IDS_PATH),
        ("evidence", PROBLEM_EVIDENCE_PATH),
        ("scores",   PROBLEM_SCORES_PATH),
        ("posts",    RAW_POSTS_PATH),
        ("evaluations", IDEA_EVALUATION_PATH)
    ]:
        if os.path.exists(path):
            try:
                dfs[name] = pd.read_csv(path, encoding="utf-8")
            except Exception:
                dfs[name] = pd.DataFrame()
        else:
            dfs[name] = pd.DataFrame()
    return dfs


def load_reviews():
    if os.path.exists(IDEA_REVIEWS_PATH):
        try:
            df = pd.read_csv(IDEA_REVIEWS_PATH)
            for col in IDEA_REVIEWS_COLUMNS:
                if col not in df.columns:
                    df[col] = ""
            return df
        except Exception:
            pass
    return pd.DataFrame(columns=IDEA_REVIEWS_COLUMNS)

def save_review(problem_id, status, reason="", notes=""):
    df = load_reviews()
    df = df[df["problem_id"] != problem_id]  # remove old entry for this problem
    new_row = pd.DataFrame([{
        "problem_id": problem_id,
        "status": status,
        "reason": reason,
        "notes": notes,
        "reviewed_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    }])
    df = pd.concat([df, new_row], ignore_index=True)
    df.to_csv(IDEA_REVIEWS_PATH, index=False)


def load_notes():
    """Load notes from manual reviews DataFrame to match expected notes schema."""
    reviews = load_reviews()
    if not reviews.empty and "notes" in reviews.columns:
        notes_df = reviews.rename(columns={"notes": "note"})
        return notes_df
    return pd.DataFrame(columns=["problem_id", "note"])


def clean_pdf_text(text):
    """Clean text for FPDF Helvetica: remove emojis and encode/decode to latin-1."""
    if not text:
        return ""
    # Remove emojis and high-unicode characters
    text = re.sub(r'[^\x00-\x7F\x80-\xFF]', '', text)
    # Replace common unicode quotes and dashes with latin-1 equivalents
    replacements = {
        '\u201c': '"', '\u201d': '"', '\u2018': "'", '\u2019': "'",
        '\u2013': '-', '\u2014': '-', '\u2022': '*', '\u2192': '->',
        '\u23f1': '', '\u2b50': '*', '\U0001f310': '', '\u2714': '[PASS]',
        '\u26a0': '[WARN]', '\u274c': '[FAIL]', '⏱️': '', '🟢': '',
        '🟡': '', '🟠': '', '🔴': '', '⬆': '', '🎯': '', '💰': '',
        '📝': '', '🧠': '', '📋': '', '📊': '', '💬': '',
        '✅': '', '⏸️': '', '❌': '', '🔄': '', '📌': '', '💡': ''
    }
    for orig, rep in replacements.items():
        text = text.replace(orig, rep)
    
    # Encode to latin-1, ignoring non-latin-1 characters
    try:
        return text.encode('latin-1', 'ignore').decode('latin-1')
    except Exception:
        return "".join(c for c in text if ord(c) < 256)


def clean_pptx_text(text):
    """Clean text for PPTX: remove emojis and keep standard unicode."""
    if not text:
        return ""
    replacements = {
        '⏱️': '', '🟢': '', '🟡': '', '🟠': '', '🔴': '', '⬆': '',
        '🎯': '', '💰': '', '📝': '', '🧠': '', '📋': '', '📊': '',
        '💬': '', '✅': '', '⏸️': '', '❌': '', '🔄': '', '📌': '', '💡': ''
    }
    for orig, rep in replacements.items():
        text = text.replace(orig, rep)
    return text


def generate_plotly_chart_image(problem_scores):
    """Generate the Plotly chart as PNG bytes."""
    if problem_scores.empty:
        return None
    try:
        problem_scores_sorted = problem_scores.sort_values("run_date")
        
        if len(problem_scores_sorted) == 1:
            latest = problem_scores_sorted.iloc[0]
            categories = ["Acuteness", "Clarity", "Market Size", "Competition", "Ideaspace", "Real Problem", "Tarpit Risk", "Proxies"]
            values = [float(latest.get(k,0)) for k in ["problem_acuteness", "customer_clarity", "market_size", "competition", "good_ideaspace", "real_problem", "tarpit_risk", "good_proxies"]]
            
            fig = go.Figure(data=go.Scatterpolar(
              r=values + [values[0]],
              theta=categories + [categories[0]],
              fill='toself',
              line=dict(color="#818cf8"),
              marker=dict(color="#c084fc")
            ))
            fig.update_layout(
              polar=dict(
                radialaxis=dict(visible=True, range=[0, 10], color="rgba(0,0,0,0.6)")
              ),
              paper_bgcolor="rgba(255,255,255,1)",
              font=dict(color="black", size=11),
              height=320,
              width=450,
              margin=dict(l=40, r=40, t=20, b=20)
            )
        else:
            fig = go.Figure()
            fig.add_trace(go.Scatter(
                x=problem_scores_sorted["run_date"],
                y=problem_scores_sorted["final_rank_score"].astype(float),
                mode="lines+markers",
                line=dict(color="#818cf8", width=3),
                marker=dict(size=8, color="#c084fc"),
                name="Final Rank Score"
            ))
            fig.update_layout(
                paper_bgcolor="rgba(255,255,255,1)",
                plot_bgcolor="rgba(240,240,240,0.5)",
                font=dict(color="black", size=11),
                height=320,
                width=550,
                margin=dict(l=40, r=40, t=20, b=20),
                xaxis=dict(showgrid=True, gridcolor="rgba(0,0,0,0.1)", title="Run Date", color="black"),
                yaxis=dict(showgrid=True, gridcolor="rgba(0,0,0,0.1)", title="Score /100", range=[0, 100], color="black")
            )
            
        img_bytes = fig.to_image(format="png")
        return img_bytes
    except Exception as e:
        print(f"Error generating chart image for export: {e}")
        return None


def build_export_data(problem, problem_scores, eval_data, notes_df, reviews_df, evidence_df, posts_df):
    """Collect all data for a problem into a clean dict for export."""
    problem_id   = problem["problem_id"]
    problem_name = str(problem.get("problem_name", "Unknown"))
    industry     = str(problem.get("industry", "Other"))
    evidence_count = int(problem.get("evidence_count", 0))
    final_score  = float(problem.get("latest_final_rank_score", 0) or 0)
    last_seen    = str(problem.get("last_seen_date", "N/A"))

    note = ""
    if not notes_df.empty:
        m = notes_df[notes_df["problem_id"] == problem_id]
        if not m.empty:
            note = str(m.iloc[-1].get("note", ""))

    review_status = "Unreviewed"
    review_reason = ""
    if not reviews_df.empty:
        m = reviews_df[reviews_df["problem_id"] == problem_id]
        if not m.empty:
            review_status = str(m.iloc[-1].get("status", "Unreviewed"))
            review_reason = str(m.iloc[-1].get("reason", ""))

    scores = {}
    if not problem_scores.empty:
        latest = problem_scores.sort_values("run_date", ascending=False).iloc[0]
        for k in ["problem_acuteness","customer_clarity","market_size","competition",
                  "good_ideaspace","real_problem","tarpit_risk","good_proxies"]:
            scores[k] = float(latest.get(k, 0) or 0)

    # Generate Plotly Chart static image
    chart_png = generate_plotly_chart_image(problem_scores)

    # Extract linked evidence posts from Reddit
    evidence_posts = []
    if not evidence_df.empty and "problem_id" in evidence_df.columns:
        problem_evidence = evidence_df[evidence_df["problem_id"] == problem_id].copy()
        if not problem_evidence.empty and not posts_df.empty:
            merged = problem_evidence.merge(
                posts_df, on="post_id", how="left", suffixes=("_ev", "_post")
            )
            if "upvotes" in merged.columns:
                merged = merged.sort_values("upvotes", ascending=False)
            
            for _, ev_row in merged.iterrows():
                evidence_posts.append({
                    "subreddit": str(ev_row.get("subreddit", "unknown")),
                    "upvotes": int(ev_row.get("upvotes", 0)),
                    "post_created_date": str(ev_row.get("post_created_date", ""))[:10],
                    "title": str(ev_row.get("title", "Untitled")),
                    "body": str(ev_row.get("body", "")),
                    "top_comments": str(ev_row.get("top_comments", "")),
                    "similarity_score": float(ev_row.get("similarity_score", 0) or 0),
                    "wtp_score": int(ev_row.get("wtp_score", ev_row.get("wtp_score_ev", 0)) or 0)
                })

    return {
        "problem_name": problem_name, "industry": industry,
        "evidence_count": evidence_count, "final_score": final_score,
        "last_seen": last_seen, "note": note,
        "review_status": review_status, "review_reason": review_reason,
        "scores": scores, "eval_data": eval_data,
        "chart_png": chart_png, "evidence_posts": evidence_posts
    }


def export_as_pdf(data):
    """Generate a PDF from problem data. Returns bytes."""
    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)

    # Title
    pdf.set_font("Helvetica", "B", 20)
    pdf.set_text_color(30, 30, 60)
    pdf.set_x(pdf.l_margin)
    pdf.cell(pdf.epw, 12, clean_pdf_text(data["problem_name"][:80]), ln=True)

    # Meta row
    pdf.set_font("Helvetica", "", 11)
    pdf.set_text_color(80, 80, 100)
    pdf.set_x(pdf.l_margin)
    pdf.cell(pdf.epw, 8,
        clean_pdf_text(
            f"Industry: {data['industry']}   |   Evidence Posts: {data['evidence_count']}   "
            f"|   Final Score: {data['final_score']:.1f}/100   |   Last Seen: {data['last_seen']}"
        ),
        ln=True)

    # Review status
    pdf.set_font("Helvetica", "B", 11)
    pdf.set_text_color(50, 50, 50)
    pdf.set_x(pdf.l_margin)
    pdf.cell(pdf.epw, 10, clean_pdf_text(f"Review Status: {data['review_status']}"), ln=True)
    if data["review_reason"]:
        pdf.set_font("Helvetica", "", 10)
        pdf.set_x(pdf.l_margin)
        pdf.cell(pdf.epw, 6, clean_pdf_text(f"Reason: {data['review_reason']}"), ln=True)

    pdf.ln(4)

    # Scores section
    if data["scores"]:
        pdf.set_font("Helvetica", "B", 13)
        pdf.set_text_color(30, 30, 60)
        pdf.set_x(pdf.l_margin)
        pdf.cell(pdf.epw, 10, "Score Breakdown", ln=True)
        pdf.set_font("Helvetica", "", 10)
        pdf.set_text_color(50, 50, 50)
        score_labels = {
            "problem_acuteness": "Problem Acuteness",
            "customer_clarity": "Customer Clarity",
            "market_size": "Market Size",
            "competition": "Competition",
            "good_ideaspace": "Good Ideaspace",
            "real_problem": "Real Problem",
            "tarpit_risk": "Tarpit Risk",
            "good_proxies": "Good Proxies"
        }
        for k, label in score_labels.items():
            val = data["scores"].get(k, 0)
            pdf.set_x(pdf.l_margin)
            pdf.cell(80, 7, clean_pdf_text(label))
            pdf.cell(0, 7, f"{val:.0f} / 10", ln=True)

    # Embed Chart image if available
    if data["chart_png"]:
        pdf.ln(6)
        pdf.set_font("Helvetica", "B", 13)
        pdf.set_text_color(30, 30, 60)
        pdf.set_x(pdf.l_margin)
        pdf.cell(pdf.epw, 10, "Score History Chart", ln=True)
        chart_io = io.BytesIO(data["chart_png"])
        pdf.set_x(pdf.l_margin)
        pdf.image(chart_io, w=110)

    pdf.ln(6)

    # Evaluation table section
    if data["eval_data"]:
        pdf.set_font("Helvetica", "B", 13)
        pdf.set_text_color(30, 30, 60)
        pdf.set_x(pdf.l_margin)
        pdf.cell(pdf.epw, 10, "Idea Evaluation Table", ln=True)
        for dim in EVALUATION_DIMENSIONS:
            key     = dim["key"]
            dim_val = data["eval_data"].get(key, {})
            verdict = str(dim_val.get("verdict", "WARN"))
            answer  = str(dim_val.get("answer", ""))
            color_map = {"PASS": (34,197,94), "WARN": (234,179,8), "FAIL": (239,68,68)}
            r, g, b = color_map.get(verdict, (150,150,150))
            pdf.set_font("Helvetica", "B", 10)
            pdf.set_text_color(r, g, b)
            pdf.set_x(pdf.l_margin)
            pdf.cell(pdf.epw, 7, clean_pdf_text(f"[{verdict}] {dim['category']}"), ln=True)
            pdf.set_font("Helvetica", "", 9)
            pdf.set_text_color(60, 60, 60)
            pdf.set_x(pdf.l_margin)
            pdf.multi_cell(pdf.epw, 6, clean_pdf_text(answer[:400]))
            pdf.ln(2)

    # Notes
    if data["note"]:
        pdf.ln(4)
        pdf.set_font("Helvetica", "B", 13)
        pdf.set_text_color(30, 30, 60)
        pdf.set_x(pdf.l_margin)
        pdf.cell(pdf.epw, 10, "My Notes", ln=True)
        pdf.set_font("Helvetica", "", 10)
        pdf.set_text_color(60, 60, 60)
        pdf.set_x(pdf.l_margin)
        pdf.multi_cell(pdf.epw, 6, clean_pdf_text(data["note"][:800]))

    # Reddit Evidence Posts
    if data["evidence_posts"]:
        pdf.ln(6)
        pdf.set_font("Helvetica", "B", 13)
        pdf.set_text_color(30, 30, 60)
        pdf.set_x(pdf.l_margin)
        pdf.cell(pdf.epw, 10, "Linked Reddit Evidence & Voice of Customer", ln=True)
        for post in data["evidence_posts"]:
            pdf.set_font("Helvetica", "B", 10)
            pdf.set_text_color(244, 114, 182) # pink color matching sub badge
            pdf.set_x(pdf.l_margin)
            pdf.cell(pdf.epw, 6, clean_pdf_text(f"r/{post['subreddit']}   |   Upvotes: {post['upvotes']}   |   Date: {post['post_created_date']}"), ln=True)
            pdf.set_font("Helvetica", "B", 9)
            pdf.set_text_color(50, 50, 50)
            pdf.set_x(pdf.l_margin)
            pdf.multi_cell(pdf.epw, 5, clean_pdf_text(f"Title: {post['title']}"))
            pdf.set_font("Helvetica", "", 8)
            pdf.set_text_color(100, 100, 100)
            body_preview = post['body'].replace("\n", " ").strip()
            body_preview = body_preview[:400] + "..." if len(body_preview) > 400 else body_preview
            pdf.set_x(pdf.l_margin)
            pdf.multi_cell(pdf.epw, 4, clean_pdf_text(body_preview))
            
            if post["top_comments"]:
                pdf.set_font("Helvetica", "I", 8)
                pdf.set_text_color(120, 120, 120)
                comments = post["top_comments"].split(" || ")[:2]
                pdf.set_x(pdf.l_margin)
                pdf.cell(pdf.epw, 4, "Top Comments:", ln=True)
                for c in comments:
                    pdf.set_x(pdf.l_margin)
                    pdf.multi_cell(pdf.epw, 4, clean_pdf_text(f"- {c[:200]}"))
            pdf.ln(3)

    return bytes(pdf.output())


def export_as_pptx(data):
    """Generate a PowerPoint from problem data. Returns bytes."""
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    def add_text_box(slide, text, left, top, width, height,
                     font_size=14, bold=False, color=(30,30,60), wrap=True):
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = wrap
        p  = tf.paragraphs[0]
        p.text = clean_pptx_text(text)
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(*color)

    blank_layout = prs.slide_layouts[6]

    # Slide 1 — Overview
    slide1 = prs.slides.add_slide(blank_layout)
    add_text_box(slide1, data["problem_name"], 0.5, 0.3, 12, 1.2,
                 font_size=28, bold=True, color=(30,30,80))
    add_text_box(slide1,
        f"Industry: {data['industry']}   |   Evidence: {data['evidence_count']} posts   "
        f"|   Score: {data['final_score']:.1f}/100   |   Status: {data['review_status']}",
        0.5, 1.6, 12, 0.5, font_size=13, color=(80,80,120))
    if data["note"]:
        add_text_box(slide1, f"📝 Notes:\n{data['note'][:500]}",
                     0.5, 2.3, 12, 4.5, font_size=12, color=(60,60,100))

    # Slide 2 — Scores & Plotly Chart
    if data["scores"]:
        slide2 = prs.slides.add_slide(blank_layout)
        add_text_box(slide2, "Score Breakdown & History",
                     0.5, 0.2, 12, 0.6, font_size=22, bold=True, color=(30,30,80))
        score_labels = [
            ("problem_acuteness","Problem Acuteness"),("customer_clarity","Customer Clarity"),
            ("market_size","Market Size"),("competition","Competition"),
            ("good_ideaspace","Good Ideaspace"),("real_problem","Real Problem"),
            ("tarpit_risk","Tarpit Risk"),("good_proxies","Good Proxies")
        ]
        col, row = 0, 0
        for key, label in score_labels:
            val = data["scores"].get(key, 0)
            x = 0.5 + col * 2.8
            y = 1.1 + row * 1.4
            color = (34,197,94) if val>=7 else (234,179,8) if val>=4 else (239,68,68)
            add_text_box(slide2, label, x, y, 2.6, 0.4,
                         font_size=10, bold=True, color=(60,60,80))
            add_text_box(slide2, f"{val:.0f}/10", x, y+0.35, 2.6, 0.7,
                         font_size=28, bold=True, color=color)
            col += 1
            if col == 2:
                col = 0
                row += 1

        # Embed Plotly Score Chart on the Right column
        if data["chart_png"]:
            chart_io = io.BytesIO(data["chart_png"])
            slide2.shapes.add_picture(chart_io, Inches(6.4), Inches(1.1), Inches(6.4), Inches(5.2))

    # Slide 3 — Evaluation Table
    if data["eval_data"]:
        slide3 = prs.slides.add_slide(blank_layout)
        add_text_box(slide3, "Idea Evaluation",
                     0.3, 0.1, 12, 0.5, font_size=20, bold=True, color=(30,30,80))
        for i, dim in enumerate(EVALUATION_DIMENSIONS[:5]):
            dim_val = data["eval_data"].get(dim["key"], {})
            verdict = str(dim_val.get("verdict", "WARN"))
            answer  = str(dim_val.get("answer", ""))[:250]
            vc = {"PASS":(34,197,94),"WARN":(234,179,8),"FAIL":(239,68,68)}
            col_idx = i % 2
            row_idx = i // 2
            x = 0.3 + col_idx * 6.5
            y = 0.8 + row_idx * 2.1
            add_text_box(slide3, f"[{verdict}] {dim['category']}",
                         x, y, 6.2, 0.4, font_size=11, bold=True, color=vc.get(verdict,(150,150,150)))
            add_text_box(slide3, answer, x, y+0.4, 6.2, 1.5, font_size=10, color=(60,60,80))

    # Slide 4 — Reddit Evidence & Voice of Customer
    if data["evidence_posts"]:
        slide4 = prs.slides.add_slide(blank_layout)
        add_text_box(slide4, "Reddit Evidence & Voice of Customer",
                     0.5, 0.2, 12, 0.6, font_size=22, bold=True, color=(30,30,80))
        
        # Display top 2 evidence posts side-by-side
        for idx, post in enumerate(data["evidence_posts"][:2]):
            x = 0.5 + idx * 6.2
            y = 1.0
            w = 5.8
            
            add_text_box(slide4, f"r/{post['subreddit']}  |  Upvotes: {post['upvotes']}",
                         x, y, w, 0.4, font_size=12, bold=True, color=(244, 114, 182))
            
            title_text = post['title']
            if len(title_text) > 80:
                title_text = title_text[:77] + "..."
            add_text_box(slide4, title_text,
                         x, y+0.4, w, 0.6, font_size=11, bold=True, color=(30,30,60))
            
            body_text = post['body'].replace("\n", " ").strip()
            if len(body_text) > 350:
                body_text = body_text[:347] + "..."
            add_text_box(slide4, body_text,
                         x, y+1.0, w, 2.0, font_size=10, color=(100,100,100))
            
            if post["top_comments"]:
                comments = post["top_comments"].split(" || ")[:1]
                if comments:
                    add_text_box(slide4, f"💬 Top Comment: \"{comments[0][:200]}\"",
                                 x, y+3.1, w, 1.8, font_size=9, color=(80,80,100))

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()



@st.cache_data(ttl=60)
def compute_live_wtp(evidence_df: pd.DataFrame, posts_df: pd.DataFrame) -> pd.DataFrame:
    """
    Re-compute wtp_score per evidence row from raw post text.
    Returns enriched evidence_df with accurate wtp_score & urgency_keywords.
    Called live because step4 is no longer executed in main.py.
    """
    if evidence_df.empty or posts_df.empty:
        return evidence_df

    ev = evidence_df.copy()
    ev["wtp_score"] = 0
    ev["urgency_keywords"] = ""

    post_lookup = posts_df.set_index("post_id") if "post_id" in posts_df.columns else pd.DataFrame()

    for idx, row in ev.iterrows():
        post_id = str(row["post_id"])
        if post_id not in post_lookup.index:
            continue
        post = post_lookup.loc[post_id]
        full_text = " ".join([
            str(post.get("title", "")),
            str(post.get("body", "")),
            str(post.get("top_comments", ""))
        ]).lower()

        wtp_count = sum(1 for p in WTP_PHRASES if p in full_text)
        ev.loc[idx, "wtp_score"] = min(wtp_count, 3)
        matched = [p for p in URGENCY_PHRASES if p in full_text]
        ev.loc[idx, "urgency_keywords"] = ", ".join(matched)

    return ev


@st.cache_data(ttl=60)
def compute_avg_wtp(problems_df: pd.DataFrame, evidence_df: pd.DataFrame) -> pd.DataFrame:
    """Recalculate avg_wtp_score per problem from live evidence WTP scores."""
    df = problems_df.copy()
    if evidence_df.empty or "wtp_score" not in evidence_df.columns:
        return df
    grouped = (
        evidence_df.groupby("problem_id")["wtp_score"]
        .apply(lambda s: round(pd.to_numeric(s, errors="coerce").fillna(0).mean(), 2))
        .reset_index()
        .rename(columns={"wtp_score": "avg_wtp_score"})
    )
    df = df.drop(columns=["avg_wtp_score"], errors="ignore")
    df = df.merge(grouped, on="problem_id", how="left")
    df["avg_wtp_score"] = df["avg_wtp_score"].fillna(0.0)
    return df


# ── Custom CSS ───────────────────────────────────────────────────
st.markdown("""
<style>
    /* Hide the Deploy toolbar — keep sidebar toggle working */
    [data-testid="stToolbar"]    { display: none !important; }
    [data-testid="stDecoration"] { display: none !important; }
    #MainMenu                    { visibility: hidden; }
    footer                       { visibility: hidden; }

    /* Force sidebar to always be visible */
    section[data-testid="stSidebar"] {
        display: block !important;
        width: 21rem !important;
        min-width: 21rem !important;
        transform: none !important;
        position: relative !important;
        visibility: visible !important;
        opacity: 1 !important;
    }
    section[data-testid="stSidebar"] > div {
        width: 21rem !important;
    }
    
    /* Hide both toggle buttons — sidebar is always open */
    [data-testid="stSidebarCollapseButton"],
    [data-testid="collapsedControl"] {
        display: none !important;
    }

    /* Style the sidebar refresh button beautifully and full-width */
    div[data-testid="stSidebar"] div[data-testid="stButton"] button {
        width: 100% !important;
        background: linear-gradient(90deg, #818cf8, #c084fc) !important;
        color: white !important;
        font-weight: 700 !important;
        font-size: 0.95rem !important;
        padding: 10px 16px !important;
        border-radius: 10px !important;
        border: none !important;
        box-shadow: 0 4px 15px rgba(129, 140, 248, 0.2) !important;
        transition: all 0.3s ease !important;
        display: flex !important;
        justify-content: center !important;
        align-items: center !important;
        gap: 8px !important;
        cursor: pointer !important;
    }
    div[data-testid="stSidebar"] div[data-testid="stButton"] button:hover {
        background: linear-gradient(90deg, #9333ea, #4f46e5) !important;
        box-shadow: 0 6px 20px rgba(129, 140, 248, 0.4) !important;
        transform: translateY(-2px) !important;
    }
    div[data-testid="stSidebar"] div[data-testid="stButton"] button:active {
        transform: translateY(1px) !important;
    }

    /* Main background */
    .stApp {
        background: linear-gradient(135deg, #0f0c29, #302b63, #24243e);
    }

    /* Sidebar */
    section[data-testid="stSidebar"] {
        background: linear-gradient(180deg, #1a1a2e, #16213e);
        border-right: 1px solid rgba(255,255,255,0.05);
    }

    /* ── Custom top header bar ─────────────────────────── */
    .top-header {
        display: flex;
        align-items: center;
        gap: 14px;
        padding: 14px 24px;
        margin-bottom: 20px;
        background: linear-gradient(90deg,
            rgba(129,140,248,0.12), rgba(192,132,252,0.08), rgba(244,114,182,0.06));
        border-bottom: 1px solid rgba(129,140,248,0.2);
        border-radius: 0 0 16px 16px;
    }
    .top-header h1 {
        font-size: 1.6rem;
        font-weight: 900;
        background: linear-gradient(90deg, #818cf8, #c084fc, #f472b6);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        margin: 0;
        letter-spacing: -0.3px;
    }
    .top-header .subtitle {
        font-size: 0.8rem;
        color: rgba(255,255,255,0.35);
        margin: 0;
        margin-left: auto;
    }

    /* Metric cards — fixed equal height */
    .metric-card {
        background: linear-gradient(135deg, rgba(255,255,255,0.06), rgba(255,255,255,0.02));
        backdrop-filter: blur(10px);
        border: 1px solid rgba(255,255,255,0.1);
        border-radius: 16px;
        padding: 20px 16px;
        text-align: center;
        height: 110px;
        box-sizing: border-box;
        display: flex;
        flex-direction: column;
        align-items: center;
        justify-content: center;
        transition: transform 0.2s ease, box-shadow 0.2s ease;
    }
    .metric-card:hover {
        transform: translateY(-2px);
        box-shadow: 0 8px 32px rgba(99, 102, 241, 0.15);
    }
    .metric-value {
        font-size: 2.2rem;
        font-weight: 700;
        background: linear-gradient(90deg, #818cf8, #c084fc);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        line-height: 1.1;
    }
    .metric-sub {
        font-size: 0.72rem;
        color: rgba(248,113,113,0.75);
        font-weight: 500;
        margin: 2px 0;
    }
    .metric-label {
        font-size: 0.75rem;
        color: rgba(255,255,255,0.45);
        text-transform: uppercase;
        letter-spacing: 1px;
        margin-top: 3px;
    }

    /* Problem row */
    .problem-row {
        background: rgba(255,255,255,0.03);
        border: 1px solid rgba(255,255,255,0.06);
        border-radius: 12px;
        padding: 14px 20px;
        margin-bottom: 8px;
        transition: background 0.2s ease;
    }
    .problem-row:hover { background: rgba(255,255,255,0.06); }

    /* Industry badge */
    .industry-badge {
        display: inline-block;
        padding: 3px 10px;
        border-radius: 20px;
        font-size: 0.75rem;
        font-weight: 600;
        background: rgba(99, 102, 241, 0.15);
        color: #818cf8;
        border: 1px solid rgba(99, 102, 241, 0.3);
    }

    /* WTP badge colors */
    .wtp-green { color: #34d399; }
    .wtp-amber { color: #fbbf24; }
    .wtp-red   { color: #f87171; }

    /* Score badge */
    .score-badge {
        display: inline-block;
        padding: 4px 12px;
        border-radius: 8px;
        font-weight: 700;
        font-size: 0.95rem;
    }
    .score-high { background: rgba(52,211,153,0.15); color:#34d399; border:1px solid rgba(52,211,153,0.3); }
    .score-mid  { background: rgba(251,191,36,0.15);  color:#fbbf24; border:1px solid rgba(251,191,36,0.3); }
    .score-low  { background: rgba(248,113,113,0.15); color:#f87171; border:1px solid rgba(248,113,113,0.3); }

    /* Evidence post card */
    .evidence-card {
        background: rgba(255,255,255,0.03);
        border: 1px solid rgba(255,255,255,0.08);
        border-radius: 10px;
        padding: 16px;
        margin-bottom: 12px;
    }
    .evidence-body {
        font-size: 0.85rem;
        color: rgba(255,255,255,0.55);
        line-height: 1.5;
        display: -webkit-box;
        -webkit-line-clamp: 2;
        -webkit-box-orient: vertical;
        overflow: hidden;
    }

    /* Subreddit badge */
    .sub-badge {
        display: inline-block;
        padding: 2px 8px;
        border-radius: 12px;
        font-size: 0.7rem;
        background: rgba(244,114,182,0.12);
        color: #f472b6;
        border: 1px solid rgba(244,114,182,0.25);
        margin-right: 6px;
    }

    /* Star styling */
    .wtp-stars { font-size: 1.1rem; }
</style>
""", unsafe_allow_html=True)


# ── Load & enrich data ───────────────────────────────────────────
data = load_data()
problems_df = data["problems"]
evidence_df = data["evidence"]
scores_df   = data["scores"]
posts_df    = data["posts"]
evaluations_df = data.get("evaluations", pd.DataFrame())

if not problems_df.empty and "last_run_timestamp" not in problems_df.columns:
    problems_df["last_run_timestamp"] = ""

# Backfill last_run_timestamp from problem_scores run_date for problems that haven't
# been run since the column was introduced. Use latest run_date per problem_id.
if not problems_df.empty and not scores_df.empty and "run_date" in scores_df.columns:
    try:
        # Get the latest run_date per problem
        latest_score_dates = (
            scores_df.dropna(subset=["run_date"])
            .sort_values("run_date")
            .groupby("problem_id")["run_date"]
            .last()
            .reset_index()
            .rename(columns={"run_date": "_latest_score_date"})
        )
        problems_df = problems_df.merge(latest_score_dates, on="problem_id", how="left")

        # Only fill in last_run_timestamp where it's missing/blank
        missing_ts = (
            problems_df["last_run_timestamp"].isna() |
            (problems_df["last_run_timestamp"].astype(str).str.strip() == "") |
            (problems_df["last_run_timestamp"].astype(str).str.strip().str.lower() == "nan")
        )
        problems_df.loc[missing_ts, "last_run_timestamp"] = (
            problems_df.loc[missing_ts, "_latest_score_date"]
        )
        problems_df = problems_df.drop(columns=["_latest_score_date"], errors="ignore")
    except Exception:
        pass

# Re-compute WTP live (step4 is no longer run in the pipeline)
if not evidence_df.empty and not posts_df.empty:
    evidence_df = compute_live_wtp(evidence_df, posts_df)
    problems_df = compute_avg_wtp(problems_df, evidence_df)


# ── Helper functions ─────────────────────────────────────────────
def wtp_color_class(score):
    if score >= 2.0: return "wtp-green"
    if score >= 1.0: return "wtp-amber"
    return "wtp-red"

def score_class(score):
    if score >= 70: return "score-high"
    if score >= 40: return "score-mid"
    return "score-low"

def wtp_stars(score):
    full  = int(score)
    half  = 1 if score - full >= 0.5 else 0
    empty = 3 - full - half
    return "★" * full + "☆" * (half + empty)


# ── Custom top header (replaces Streamlit's Deploy bar) ──────────
st.markdown("""
<div class="top-header">
    <span style="font-size:1.8rem;">🔍</span>
    <h1>Problem Discovery Dashboard</h1>
    <p class="subtitle">Startup problems from Reddit · Scored by AI for market potential</p>
</div>
""", unsafe_allow_html=True)


# ── Top Metric Cards ─────────────────────────────────────────────
total_problems = len(problems_df) if not problems_df.empty else 0
scanned_count  = len(posts_df)   if not posts_df.empty    else 0
noise_count    = 0
if not posts_df.empty and "passed_noise_filter" in posts_df.columns:
    noise_count = int(posts_df["passed_noise_filter"].astype(str).str.lower().isin(["false","0","0.0"]).sum())

last_run = "Never"
if not posts_df.empty and "scraped_at" in posts_df.columns:
    try:
        last_run = pd.to_datetime(posts_df["scraped_at"]).max().strftime("%Y-%m-%d %H:%M")
    except Exception:
        last_run = "Unknown"

avg_rank = f"{problems_df['latest_final_rank_score'].mean():.1f}" if not problems_df.empty and 'latest_final_rank_score' in problems_df.columns else "0.0"

col1, col2, col3, col4 = st.columns(4)
with col1:
    st.markdown(f"""
    <div class="metric-card">
        <div class="metric-value">{total_problems}</div>
        <div class="metric-label">Total Problems Tracked</div>
    </div>""", unsafe_allow_html=True)

with col2:
    st.markdown(f"""
    <div class="metric-card">
        <div class="metric-value">{scanned_count}</div>
        <div class="metric-sub">🗑️ {noise_count} filtered as noise</div>
        <div class="metric-label">Scanned Posts</div>
    </div>""", unsafe_allow_html=True)

with col3:
    st.markdown(f"""
    <div class="metric-card">
        <div class="metric-value" style="font-size:1.3rem;">{last_run}</div>
        <div class="metric-label">Last Run Date</div>
    </div>""", unsafe_allow_html=True)

with col4:
    st.markdown(f"""
    <div class="metric-card">
        <div class="metric-value">{avg_rank}</div>
        <div class="metric-label">Avg Rank Score</div>
    </div>""", unsafe_allow_html=True)

st.markdown("<br>", unsafe_allow_html=True)


# ── Sidebar Filters ───────────────────────────────────────────────
with st.sidebar:
    st.markdown("<h3 style='text-align: center; color: white; margin-top: 10px; margin-bottom: 15px; font-weight: 700; letter-spacing: 0.5px;'>🎯 Filters & Controls</h3>", unsafe_allow_html=True)
    
    # Premium Centered Full-Width Refresh Button
    if st.button("🔄 Refresh Dashboard", use_container_width=True, help="Clear cache and reload raw data"):
        st.cache_data.clear()
        st.rerun()
        
    st.markdown("<hr style='margin: 20px 0; border: none; border-top: 1px dashed rgba(255,255,255,0.15);'>", unsafe_allow_html=True)

    # Industry filter
    industries = ["All"]
    if not problems_df.empty and "industry" in problems_df.columns:
        industries += sorted(problems_df["industry"].dropna().unique().tolist())
    selected_industry = st.selectbox("Industry", industries, index=0)

    # Pipeline Run Timestamp filter
    run_timestamps = ["All"]
    if not problems_df.empty and "last_run_timestamp" in problems_df.columns:
        raw_runs = problems_df["last_run_timestamp"].dropna().unique().tolist()
        sorted_runs = sorted([str(r) for r in raw_runs if str(r).strip() and str(r).strip().lower() != "nan"], reverse=True)
        if sorted_runs:
            run_timestamps += ["Latest Run"] + sorted_runs
            
    default_run_index = 0
    if len(run_timestamps) > 1 and "Latest Run" in run_timestamps:
        default_run_index = 1  # Default to Latest Run
        
    selected_run = st.selectbox("Pipeline Run (Completion Time)", run_timestamps, index=default_run_index)

    # Min evidence count
    min_evidence = st.slider("Min Evidence Count", 1, 20, 1)

    # Min WTP score
    min_wtp = st.slider("Min WTP Score", 0.0, 3.0, 0.0, step=0.5)

    # Min LLM score (latest_total_score, max 80)
    min_llm = st.slider("Min LLM Score (/80)", 0, 80, 0, step=5)

    # Min Final Rank score (latest_final_rank_score, max 100)
    min_final = st.slider("Min Final Score (/100)", 0, 100, 0, step=5)

    # Sort by
    sort_options = {
        "Final Rank Score": "latest_final_rank_score",
        "WTP Score":        "avg_wtp_score",
        "Evidence Count":   "evidence_count",
        "Last Seen Date":   "last_seen_date",
        "Last Run Timestamp": "last_run_timestamp"
    }
    sort_by_label = st.selectbox("Sort By", list(sort_options.keys()))
    sort_by = sort_options[sort_by_label]

    # Top N
    top_n = int(st.number_input("Show Top N", min_value=1, value=10, step=1))

    st.markdown("---")
    st.markdown("### 🗂️ Review Filter")
    review_filter = st.multiselect(
        "Show by Review Status",
        options=["Unreviewed", "Accepted", "Hold", "Rejected"],
        default=["Unreviewed", "Accepted", "Hold"]
    )
    
    # st.markdown("---")🎯 Filters

    # st.markdown("### 📊 Database Summary")
    # sub_count = len(posts_df['subreddit'].unique()) if not posts_df.empty and 'subreddit' in posts_df.columns else 0
    # st.markdown(f"- **{sub_count}** Subreddits Tracking")
    # st.markdown(f"- **{len(WTP_PHRASES)}** WTP Phrases")
    # st.markdown(f"- **{len(URGENCY_PHRASES)}** Urgency Phrases")
    # st.markdown("- **Model:** Llama 3.3 70B")


# ── Apply Filters ────────────────────────────────────────────────
if problems_df.empty:
    st.info(
        "📭 **No data yet.** Run the pipeline first:\n\n"
        "```\npython main.py\n```\n\nThen refresh this page."
    )
    st.stop()

filtered_df = problems_df.copy()

# Convert numeric columns
for col in ["evidence_count", "avg_wtp_score", "latest_total_score", "latest_final_rank_score"]:
    if col in filtered_df.columns:
        filtered_df[col] = pd.to_numeric(filtered_df[col], errors="coerce").fillna(0)

# Apply filters
if selected_industry != "All":
    filtered_df = filtered_df[filtered_df["industry"] == selected_industry]

if "last_run_timestamp" in filtered_df.columns:
    if selected_run == "Latest Run":
        all_timestamps = filtered_df["last_run_timestamp"].dropna().unique()
        valid_timestamps = [t for t in all_timestamps if str(t).strip() and str(t).strip().lower() != "nan"]
        if valid_timestamps:
            latest_time = max(valid_timestamps)
            filtered_df = filtered_df[filtered_df["last_run_timestamp"] == latest_time]
    elif selected_run != "All":
        filtered_df = filtered_df[filtered_df["last_run_timestamp"] == selected_run]

filtered_df = filtered_df[filtered_df["evidence_count"]           >= min_evidence]
filtered_df = filtered_df[filtered_df["avg_wtp_score"]            >= min_wtp]
filtered_df = filtered_df[filtered_df["latest_total_score"]       >= min_llm]
filtered_df = filtered_df[filtered_df["latest_final_rank_score"]  >= min_final]

# Apply review filter
reviews_df = load_reviews()

def get_review_status(problem_id):
    if reviews_df.empty:
        return "Unreviewed"
    match = reviews_df[reviews_df["problem_id"] == problem_id]
    if match.empty:
        return "Unreviewed"
    return str(match.iloc[-1]["status"])

filtered_df["_review_status"] = filtered_df["problem_id"].apply(get_review_status)
filtered_df = filtered_df[filtered_df["_review_status"].isin(review_filter)]

# Sort
if sort_by in filtered_df.columns:
    filtered_df = filtered_df.sort_values(sort_by, ascending=False)

# Limit
filtered_df = filtered_df.head(top_n)

if filtered_df.empty:
    st.warning("No problems match the current filters. Try adjusting the sidebar filters.")
    st.stop()


# ── Main Problem List ────────────────────────────────────────────
st.markdown(f"### Showing {len(filtered_df)} problems")

for rank, (idx, problem) in enumerate(filtered_df.iterrows(), 1):
    problem_id    = problem["problem_id"]
    problem_name  = html_lib.escape(str(problem.get("problem_name", "Unknown")))
    industry      = problem.get("industry",               "Other")
    evidence_count = int(problem.get("evidence_count",    0))
    avg_wtp        = float(problem.get("avg_wtp_score",   0))
    final_score    = float(problem.get("latest_final_rank_score", 0))
    total_score    = float(problem.get("latest_total_score",      0))
    last_seen      = str(problem.get("last_seen_date",    "N/A"))

    # Freshness calculation
    freshness = "Old"
    if last_seen != "N/A":
        try:
            days_old = (datetime.now().date() - datetime.strptime(last_seen[:10], "%Y-%m-%d").date()).days
            if days_old <= 7: freshness = "🟢 Fresh"
            elif days_old <= 14: freshness = "🟡 Recent"
            elif days_old <= 30: freshness = "🟠 Aging"
            else: freshness = "🔴 Old"
        except ValueError:
            pass

    wtp_cls   = wtp_color_class(avg_wtp)
    score_cls = score_class(final_score)

    # UI 3: Hide 0.0 badge if unscored
    score_badge_html = f'<span class="score-badge {score_cls}">{final_score:.1f}/100</span>' if final_score > 0 else '<span style="font-size:0.8rem; color:rgba(255,255,255,0.4);">No Score</span>'

    last_run_time = str(problem.get("last_run_timestamp", ""))
    last_run_display = f'<span style="font-size:0.75rem; padding: 2px 6px; border-radius: 4px; background: rgba(99,102,241,0.15); color:#818cf8; border:1px solid rgba(99,102,241,0.25);">⏱️ Run: {last_run_time}</span>' if last_run_time and last_run_time.strip() and last_run_time.strip().lower() != "nan" else ""

    header_html = f"""<div class="problem-row">
<div style="display:flex; align-items:center; justify-content:space-between; flex-wrap:wrap; gap:8px;">
<div style="display:flex; align-items:center; gap:12px;">
<span style="font-size:1.3rem; font-weight:800; color:rgba(255,255,255,0.3);">#{rank}</span>
<span style="font-size:1.05rem; font-weight:600; color:#e2e8f0;">{problem_name}</span>
<span class="industry-badge">{industry}</span>
<span style="font-size:0.75rem; padding: 2px 6px; border-radius: 4px; background: rgba(255,255,255,0.1);">{freshness}</span>
{last_run_display}
</div>
<div style="display:flex; align-items:center; gap:16px; flex-wrap:wrap;">
<span style="font-size:0.85rem; color:rgba(255,255,255,0.5);">📝 {evidence_count}</span>
<span class="wtp-stars {wtp_cls}">WTP: {avg_wtp:.1f}</span>
{f'<span style="font-size:0.8rem; font-weight:600; color:#a78bfa; background:rgba(167,139,250,0.12); border:1px solid rgba(167,139,250,0.25); border-radius:6px; padding:2px 8px;">LLM: {total_score:.0f}/80</span>' if total_score > 0 else ''}
{score_badge_html}
</div>
</div>
</div>"""
    st.markdown(header_html, unsafe_allow_html=True)

    # Expandable detail
    with st.expander(f"📊 Details — {problem_name}", expanded=(len(filtered_df) == 1)):

        # ── Load data for export & widgets ────────────────────────
        problem_scores_exp = pd.DataFrame()
        if not scores_df.empty and "problem_id" in scores_df.columns:
            problem_scores_exp = scores_df[scores_df["problem_id"] == problem_id].copy()

        problem_eval_exp = pd.DataFrame()
        if not evaluations_df.empty and "problem_id" in evaluations_df.columns:
            problem_eval_exp = evaluations_df[evaluations_df["problem_id"] == problem_id].copy()

        # Build export data
        notes_df_exp   = load_notes()
        reviews_df_exp = load_reviews()
        eval_data_exp  = {}
        if not problem_eval_exp.empty:
            try:
                eval_data_exp = json.loads(
                    str(problem_eval_exp.sort_values("run_date", ascending=False)
                        .iloc[0].get("evaluation_json", "{}"))
                )
            except Exception:
                eval_data_exp = {}

        export_data = build_export_data(
            problem, problem_scores_exp, eval_data_exp, notes_df_exp, reviews_df_exp, evidence_df, posts_df
        )

        # ── Premium Header Action Bar ──
        action_col1, action_col2, action_col3 = st.columns([5.5, 2.25, 2.25])
        with action_col1:
            st.markdown(f"<div style='padding-top:6px; font-size:0.85rem; font-weight:800; color:#a5b4fc; letter-spacing:0.8px;'>📋 PROBLEM DISCOVERY & AI ANALYSIS REPORT</div>", unsafe_allow_html=True)
        with action_col2:
            pdf_bytes = b""
            try:
                pdf_bytes = export_as_pdf(export_data)
            except Exception as e:
                import traceback
                print(f"Error generating PDF for {problem_name}: {e}")
                traceback.print_exc()
            
            if pdf_bytes:
                st.download_button(
                    label="📄 PDF Report",
                    data=pdf_bytes,
                    file_name=f"{problem_name[:40].replace(' ','_')}_evaluation.pdf",
                    mime="application/pdf",
                    key=f"pdf_{problem_id}",
                    use_container_width=True
                )
            else:
                st.button("📄 PDF (Failed to generate)", key=f"pdf_failed_{problem_id}", disabled=True, use_container_width=True)
                
        with action_col3:
            pptx_bytes = b""
            try:
                pptx_bytes = export_as_pptx(export_data)
            except Exception as e:
                import traceback
                print(f"Error generating PPTX for {problem_name}: {e}")
                traceback.print_exc()
                
            if pptx_bytes:
                st.download_button(
                    label="📊 PPTX Deck",
                    data=pptx_bytes,
                    file_name=f"{problem_name[:40].replace(' ','_')}_evaluation.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key=f"pptx_{problem_id}",
                    use_container_width=True
                )
            else:
                st.button("📊 PPTX (Failed to generate)", key=f"pptx_failed_{problem_id}", disabled=True, use_container_width=True)
        st.markdown("<hr style='margin: 12px 0 16px 0; border: none; border-top: 1px solid rgba(255,255,255,0.08);'>", unsafe_allow_html=True)

        # ── Manual Review ─────────────────────────────────────────
        st.markdown("##### ✅ Manual Review")

        current_status = get_review_status(problem_id)
        current_reason = ""
        current_notes = ""
        if not reviews_df.empty:
            match = reviews_df[reviews_df["problem_id"] == problem_id]
            if not match.empty:
                val_r = match.iloc[-1].get("reason", "")
                if pd.notna(val_r) and str(val_r).strip().lower() != "nan":
                    current_reason = str(val_r)
                val_n = match.iloc[-1].get("notes", "")
                if pd.notna(val_n) and str(val_n).strip().lower() != "nan":
                    current_notes = str(val_n)

        status_colors = {
            "Accepted": "#34d399", "Hold": "#fbbf24",
            "Rejected": "#f87171", "Unreviewed": "#94a3b8"
        }
        status_color = status_colors.get(current_status, "#94a3b8")

        rev_col1, rev_col2, rev_col3 = st.columns([1.5, 2.5, 3.8])

        with rev_col1:
            st.markdown("<div style='padding-top: 5px;'>", unsafe_allow_html=True)
            st.markdown(
                f'<span style="font-size:0.85rem; font-weight:700; color:{status_color}; '
                f'background:rgba(255,255,255,0.06); border:1px solid {status_color}44; '
                f'border-radius:20px; padding:6px 14px; display:inline-block;">Current: {current_status}</span>',
                unsafe_allow_html=True
            )
            if current_reason:
                st.markdown(f"<div style='font-size:0.75rem; color:rgba(255,255,255,0.5); margin-top:8px;'>📝 Reason: {current_reason}</div>", unsafe_allow_html=True)
            if current_notes:
                st.markdown(f"<div style='font-size:0.75rem; color:rgba(255,255,255,0.4); margin-top:4px;'>📌 Notes: {current_notes}</div>", unsafe_allow_html=True)
            st.markdown("</div>", unsafe_allow_html=True)

        with rev_col2:
            reason_input = st.text_input(
                "Reason (optional)",
                value=current_reason,
                key=f"reason_{problem_id}",
                placeholder="Reason (e.g. Market size...)",
                label_visibility="collapsed"
            )
            notes_input = st.text_input(
                "Notes (optional)",
                value=current_notes,
                key=f"notes_{problem_id}",
                placeholder="Additional notes / comments...",
                label_visibility="collapsed"
            )

        with rev_col3:
            st.markdown("<div style='padding-top: 15px;'>", unsafe_allow_html=True)
            btn_sub1, btn_sub2, btn_sub3, btn_sub4 = st.columns([1, 1, 1, 1.2])
            with btn_sub1:
                if st.button("✅ Accept", key=f"accept_{problem_id}"):
                    save_review(problem_id, "Accepted", reason_input, notes_input)
                    st.cache_data.clear()
                    st.rerun()
            with btn_sub2:
                if st.button("⏸️ Hold", key=f"hold_{problem_id}"):
                    save_review(problem_id, "Hold", reason_input, notes_input)
                    st.cache_data.clear()
                    st.rerun()
            with btn_sub3:
                if st.button("❌ Reject", key=f"reject_{problem_id}"):
                    save_review(problem_id, "Rejected", reason_input, notes_input)
                    st.cache_data.clear()
                    st.rerun()
            with btn_sub4:
                if st.button("🔄 Reset", key=f"reset_{problem_id}", help="Remove review vote and set to Unreviewed"):
                    df = load_reviews()
                    df = df[df["problem_id"] != problem_id]
                    df.to_csv(IDEA_REVIEWS_PATH, index=False)
                    st.cache_data.clear()
                    st.rerun()
            st.markdown("</div>", unsafe_allow_html=True)

        st.markdown("---")

        # Load scores for this problem (shared by both columns)
        problem_scores = pd.DataFrame()
        if not scores_df.empty and "problem_id" in scores_df.columns:
            problem_scores = scores_df[scores_df["problem_id"] == problem_id].copy()

        detail_col1, detail_col2 = st.columns([1, 1])

        # ── Score Breakdown ──────────────────────────────────────
        with detail_col1:
            st.markdown("##### 📈 Score Breakdown")

            if not problem_scores.empty:
                latest = problem_scores.sort_values("run_date", ascending=False).iloc[0]

                factor_names = [
                    "Problem Acuteness", "Customer Clarity", "Market Size",
                    "Competition", "Good Ideaspace", "Real Problem",
                    "Tarpit Risk", "Good Proxies"
                ]
                factor_keys = [
                    "problem_acuteness", "customer_clarity", "market_size",
                    "competition", "good_ideaspace", "real_problem",
                    "tarpit_risk", "good_proxies"
                ]

                row1_cols = st.columns(4)
                row2_cols = st.columns(4)
                all_cols  = row1_cols + row2_cols
                
                # Add WTP Tile above factors
                # st.markdown(f'<div style="margin-bottom:12px; padding:10px; background:rgba(251,191,36,0.1); border:1px solid rgba(251,191,36,0.3); border-radius:8px; display:inline-block;"><div style="font-size:0.7rem; color:#fbbf24; text-transform:uppercase;">Willingness to Pay</div><div style="font-size:1.3rem; font-weight:bold; color:#fbbf24;">{avg_wtp:.1f}/3.0</div></div>', unsafe_allow_html=True)

                for col_widget, name, key in zip(all_cols, factor_names, factor_keys):
                    val = float(latest.get(key, 0))
                    color = "#34d399" if val >= 7 else ("#fbbf24" if val >= 4 else "#f87171")
                    with col_widget:
                        st.markdown(
                            f'<div style="text-align:center; padding:8px 4px; '
                            f'background:rgba(255,255,255,0.03); border:1px solid rgba(255,255,255,0.08); '
                            f'border-radius:8px;">'
                            f'<div style="font-size:0.7rem; color:rgba(255,255,255,0.45); '
                            f'text-transform:uppercase; letter-spacing:0.5px; margin-bottom:4px;">{name}</div>'
                            f'<div style="font-size:1.15rem; font-weight:800; color:{color};">{val:.0f}'
                            f'<span style="font-size:0.75rem; color:rgba(255,255,255,0.25);">/10</span></div>'
                            f'</div>',
                            unsafe_allow_html=True
                        )
            else:
                st.info("No scores yet — run step5_groq_score.py to score this problem")
                st.markdown(f"""
                <div style="display:flex; gap:20px; margin-top:10px;">
                    <div style="text-align:center; padding:12px; background:rgba(255,255,255,0.05); border-radius:8px; min-width:100px;">
                        <div style="font-size:2rem; font-weight:800;">{evidence_count}</div>
                        <div style="font-size:0.8rem; color:gray; text-transform:uppercase;">Evidence</div>
                    </div>
                    <div style="text-align:center; padding:12px; background:rgba(251,191,36,0.1); border:1px solid rgba(251,191,36,0.3); border-radius:8px; min-width:100px;">
                        <div style="font-size:2rem; font-weight:800; color:#fbbf24;">{avg_wtp:.1f}</div>
                        <div style="font-size:0.8rem; color:#fbbf24; text-transform:uppercase;">WTP Score</div>
                    </div>
                </div>
                """, unsafe_allow_html=True)

        # ── Score History Chart ──────────────────────────────────
        with detail_col2:
            st.markdown("##### 📉 Score History")

            if not problem_scores.empty:
                problem_scores_sorted = problem_scores.sort_values("run_date")
                
                if len(problem_scores_sorted) == 1:
                    latest = problem_scores_sorted.iloc[0]
                    categories = ["Acuteness", "Clarity", "Market Size", "Competition", "Ideaspace", "Real Problem", "Tarpit Risk", "Proxies"]
                    values = [float(latest.get(k,0)) for k in ["problem_acuteness", "customer_clarity", "market_size", "competition", "good_ideaspace", "real_problem", "tarpit_risk", "good_proxies"]]
                    
                    fig2 = go.Figure(data=go.Scatterpolar(
                      r=values + [values[0]],
                      theta=categories + [categories[0]],
                      fill='toself',
                      line=dict(color="#818cf8"),
                      marker=dict(color="#c084fc")
                    ))
                    fig2.update_layout(
                      polar=dict(
                        radialaxis=dict(visible=True, range=[0, 10], color="rgba(255,255,255,0.5)")
                      ),
                      paper_bgcolor="rgba(0,0,0,0)",
                      font=dict(color="white"),
                      height=320,
                      margin=dict(l=30, r=30, t=20, b=20)
                    )
                else:
                    fig2 = go.Figure()
                    fig2.add_trace(go.Scatter(
                        x=problem_scores_sorted["run_date"],
                        y=problem_scores_sorted["final_rank_score"].astype(float),
                        mode="lines+markers",
                        line=dict(color="#818cf8", width=3),
                        marker=dict(size=8, color="#c084fc"),
                        name="Final Rank Score"
                    ))
                    fig2.update_layout(
                        paper_bgcolor="rgba(0,0,0,0)",
                        plot_bgcolor="rgba(0,0,0,0)",
                        font=dict(color="white"),
                        height=320,
                        margin=dict(l=10, r=10, t=10, b=10),
                        xaxis=dict(showgrid=False, title="Run Date"),
                        yaxis=dict(showgrid=True, gridcolor="rgba(255,255,255,0.05)",
                                   title="Score /100", range=[0, 100])
                    )
                st.plotly_chart(fig2, use_container_width=True,
                                key=f"score_history_{problem_id}")
            else:
                st.info("No score history available yet")

        # ── Evidence Posts ────────────────────────────────────────
        with st.expander("📋 Evidence Posts", expanded=False):
            problem_evidence = pd.DataFrame()
            if not evidence_df.empty and "problem_id" in evidence_df.columns:
                problem_evidence = evidence_df[evidence_df["problem_id"] == problem_id].copy()

            if not problem_evidence.empty and not posts_df.empty:
                merged = problem_evidence.merge(
                    posts_df, on="post_id", how="left", suffixes=("_ev", "_post")
                )
                if "upvotes" in merged.columns:
                    merged = merged.sort_values("upvotes", ascending=False)

                for _, ev_row in merged.iterrows():
                    subreddit  = ev_row.get("subreddit",        "unknown")
                    upvotes    = int(ev_row.get("upvotes",      0))
                    post_date  = str(ev_row.get("post_created_date", ""))[:10]
                    title      = html_lib.escape(str(ev_row.get("title", "Untitled")))

                    body_raw = str(ev_row.get("body", ""))
                    body_raw = body_raw.replace("\n"," ").replace("\r"," ").replace("&#x200B;","").strip()
                    body_raw = html_lib.unescape(body_raw)
                    body_raw = re.sub(r" +", " ", body_raw)
                    body_text = body_raw[:400] if len(body_raw) > 400 else body_raw
                    body      = html_lib.escape(body_text)

                    top_comments   = str(ev_row.get("top_comments", ""))
                    post_url       = ev_row.get("post_url", "")
                    wtp_val        = int(ev_row.get("wtp_score", ev_row.get("wtp_score_ev", 0)) or 0)
                    urgency        = str(ev_row.get("urgency_keywords", ev_row.get("urgency_keywords_ev", "")) or "")
                    similarity_val = float(ev_row.get("similarity_score", 0) or 0)
                    stars          = "★" * wtp_val + "☆" * (3 - wtp_val)

                    urgency_html = f'<span style="font-size:0.75rem; color:rgba(255,255,255,0.35);">🔑 {html_lib.escape(urgency)}</span>' if urgency else ''
                    safe_url = html_lib.escape(str(post_url), quote=True)

                    card_html = (
                        '<div class="evidence-card">'
                        '<div style="display:flex; align-items:center; gap:10px; margin-bottom:8px;">'
                        f'<span class="sub-badge">r/{html_lib.escape(str(subreddit))}</span>'
                        f'<span style="font-size:0.85rem; font-weight:700; color:#fb923c; background:rgba(251,146,60,0.12); border:1px solid rgba(251,146,60,0.25); border-radius:6px; padding:2px 8px;">⬆ {upvotes:,}</span>'
                        f'<span style="font-size:0.8rem; color:rgba(255,255,255,0.35);">{post_date}</span>'
                        '</div>'
                        f'<div style="font-weight:600; color:#e2e8f0; margin-bottom:6px;">{title}</div>'
                        f'<div class="evidence-body">{body}</div>'
                        '<div style="display:flex; gap:12px; align-items:center; margin-top:8px; flex-wrap:wrap;">'
                        f'<span style="font-size:0.75rem; color:rgba(255,255,255,0.4);">🎯 Match: <span style="color:#a78bfa; font-weight:600;">{similarity_val:.0%}</span></span>'
                        f'<span style="font-size:0.75rem; color:rgba(255,255,255,0.4);">💰 WTP: <span style="color:#fbbf24; font-weight:600;">{wtp_val}/3</span> {stars}</span>'
                        f'{urgency_html}'
                        '</div>'
                        f'<div style="margin-top:8px;"><a href="{safe_url}" target="_blank" style="display:inline-block; margin-top:10px; padding:7px 18px; background:linear-gradient(135deg,rgba(255,69,0,0.18),rgba(255,69,0,0.08)); color:#ff6b35; border:1px solid rgba(255,69,0,0.35); border-radius:8px; font-size:0.82rem; font-weight:600; text-decoration:none; letter-spacing:0.3px;">🔗 View on Reddit ↗</a></div>'
                        '</div>'
                    )
                    st.markdown(card_html, unsafe_allow_html=True)

                    if top_comments and top_comments != "nan":
                        comments_list = top_comments.split(" || ")
                        with st.expander("💬 Top comments", expanded=False):
                            for ci, comment in enumerate(comments_list):
                                st.markdown(f"**Comment {ci+1}:** {comment}")
            else:
                st.info("No evidence posts linked to this problem")

        # ── Score History Table ───────────────────────────────────
        if not problem_scores.empty:
            with st.expander("📊 Score History Table", expanded=False):
                display_cols = [
                    "run_date", "wtp_score", "problem_acuteness", "customer_clarity", "market_size",
                    "competition", "good_ideaspace", "real_problem", "tarpit_risk",
                    "good_proxies", "total_score", "freshness_weight", "final_rank_score"
                ]
                available_cols = [c for c in display_cols if c in problem_scores.columns]
                st.dataframe(
                    problem_scores[available_cols].sort_values("run_date", ascending=False),
                    use_container_width=True,
                    hide_index=True
                )

        # ── Idea Evaluation Table ─────────────────────────────────
        with st.expander("🧠 Idea Evaluation Table", expanded=False):
            problem_eval = pd.DataFrame()
            if not evaluations_df.empty and "problem_id" in evaluations_df.columns:
                problem_eval = evaluations_df[
                    evaluations_df["problem_id"] == problem_id
                ].copy()

            if not problem_eval.empty:
                latest_eval = problem_eval.sort_values("run_date", ascending=False).iloc[0]
                st.caption(f"Last evaluated: {latest_eval.get('run_date', '')}")
                try:
                    eval_data = json.loads(str(latest_eval.get("evaluation_json", "{}")))
                except Exception:
                    eval_data = {}

                if eval_data:
                    verdict_colors = {
                        "PASS": ("#34d399", "rgba(52,211,153,0.08)", "rgba(52,211,153,0.35)"),
                        "WARN": ("#fbbf24", "rgba(251,191,36,0.08)", "rgba(251,191,36,0.35)"),
                        "FAIL": ("#f87171", "rgba(248,113,113,0.08)", "rgba(248,113,113,0.35)"),
                    }
                    rows_html = ""
                    for dim in EVALUATION_DIMENSIONS:
                        key      = dim["key"]
                        dim_data = eval_data.get(key, {})
                        verdict  = str(dim_data.get("verdict", "WARN")).upper()
                        answer   = str(dim_data.get("answer", "No analysis available."))
                        if verdict not in verdict_colors:
                            verdict = "WARN"
                        text_col, bg_col, border_col = verdict_colors[verdict]
                        rows_html += f"""
                        <tr style="background:{bg_col}; border-bottom:1px solid rgba(255,255,255,0.05);">
                          <td style="padding:10px 14px; vertical-align:top; width:22%;">
                            <div style="font-weight:700; color:#e2e8f0; font-size:0.83rem;">{dim['category']}</div>
                            <div style="color:rgba(255,255,255,0.4); font-size:0.73rem; margin-top:2px;">{dim['question']}</div>
                          </td>
                          <td style="padding:10px 14px; vertical-align:middle; text-align:center; width:10%;">
                            <span style="display:inline-block; padding:3px 12px; border-radius:20px;
                              font-weight:700; font-size:0.72rem; letter-spacing:0.5px;
                              color:{text_col}; background:rgba(0,0,0,0.2); border:1px solid {border_col};">
                              {verdict}
                            </span>
                          </td>
                          <td style="padding:10px 14px; vertical-align:top; color:rgba(255,255,255,0.78);
                            font-size:0.83rem; line-height:1.55; width:68%;">{answer}</td>
                        </tr>"""

                    st.markdown(f"""
                    <table style="width:100%; border-collapse:collapse; font-size:0.84rem;
                      border:1px solid rgba(255,255,255,0.08); border-radius:8px; overflow:hidden;">
                      <thead>
                        <tr style="background:rgba(129,140,248,0.12);">
                          <th style="padding:10px 14px; text-align:left; color:#a5b4fc;
                            font-size:0.75rem; text-transform:uppercase; letter-spacing:0.5px;">Category</th>
                          <th style="padding:10px 14px; text-align:center; color:#a5b4fc;
                            font-size:0.75rem; text-transform:uppercase; letter-spacing:0.5px;">Verdict</th>
                          <th style="padding:10px 14px; text-align:left; color:#a5b4fc;
                            font-size:0.75rem; text-transform:uppercase; letter-spacing:0.5px;">Analysis</th>
                        </tr>
                      </thead>
                      <tbody>{rows_html}</tbody>
                    </table>
                    """, unsafe_allow_html=True)
            else:
                st.info("No idea evaluation yet — run `python main.py` to generate.")

    st.markdown("")  # spacer between problems
